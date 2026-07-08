"""
slide_processor.py — Xử lý PPTX / PDF → thumbnail PNG
Dùng COM automation (PowerPoint) để render PPTX, PyMuPDF cho PDF.
"""

from __future__ import annotations

import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, List, Optional


# ═══════════════════════════════════════════════════════════════════════════
#  DATA CLASSES
# ═══════════════════════════════════════════════════════════════════════════

@dataclass
class SlideInfo:
    """Thông tin một slide đã được render."""
    index: int                  # 0-based index
    image_path: str             # đường dẫn tới file PNG thumbnail
    title: str = ""             # tiêu đề slide (nếu lấy được)
    assigned_pos: int = -1      # character offset trong kịch bản (-1 = chưa gán)
    assigned_text: str = ""     # đoạn văn đầu đoạn (preview ~40 chars)
    video_path: str = ""        # đường dẫn tới file MP4 video của slide (nếu có GIF/Video)
    gifs: list = field(default_factory=list) # Danh sách ảnh GIF nhúng trong slide: [{"gif_path": str, "x_rel": float, "y_rel": float, "w_rel": float, "h_rel": float}]

    @property
    def is_assigned(self) -> bool:
        return self.assigned_pos >= 0

    @property
    def display_number(self) -> int:
        """Số thứ tự hiển thị (1-based)."""
        return self.index + 1


@dataclass
class SlideAssignment:
    """Mapping giữa một slide và vị trí trong kịch bản."""
    slide_index: int
    text_position: int      # character offset trong văn bản
    text_snippet: str       # ~50 chars preview


# ═══════════════════════════════════════════════════════════════════════════
#  EXCEPTIONS
# ═══════════════════════════════════════════════════════════════════════════

class SlideProcessorError(Exception):
    pass


# ═══════════════════════════════════════════════════════════════════════════
#  INTERNAL HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def _temp_dir() -> str:
    """Tạo thư mục tạm riêng cho thumbnails."""
    d = tempfile.mkdtemp(prefix="kath_slides_")
    return d


def _extract_pptx_titles(pptx_path: str) -> List[str]:
    """Đọc tiêu đề từng slide qua python-pptx (không cần render)."""
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(pptx_path)
        titles: List[str] = []
        for slide in prs.slides:
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            titles.append(title)
        return titles
    except Exception:
        return []


# ═══════════════════════════════════════════════════════════════════════════
#  PPTX → THUMBNAILS  (via PowerPoint COM on Windows)
# ═══════════════════════════════════════════════════════════════════════════

def load_pptx(
    pptx_path: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
) -> List[SlideInfo]:
    """
    Render PPTX thành list SlideInfo (có hình ảnh PNG gốc sắc nét và thông tin GIF đè).
    """
    pptx_path = str(Path(pptx_path).resolve())

    def _report(pct: int, msg: str):
        if progress_cb:
            progress_cb(pct, msg)

    _report(0, "Đang mở file PPTX...")

    try:
        return _pptx_via_com_export_slides(pptx_path, _report)
    except Exception as e1:
        _report(10, f"COM export slides thất bại ({e1}), thử convert PDF...")
        try:
            return _pptx_via_com_to_pdf(pptx_path, _report)
        except Exception as e2:
            _report(20, f"COM PDF thất bại ({e2}), dùng fallback text...")
            return _pptx_fallback_text(pptx_path, _report)


def _pptx_via_com_export_video(
    pptx_path: str,
    report: Callable[[int, str], None],
) -> List[SlideInfo]:
    """Sử dụng PowerPoint COM để xuất presentation thành video MP4, rồi cắt thành từng slide."""
    import comtypes.client  # type: ignore
    import comtypes
    import time
    import subprocess
    import sys
    import shutil

    # Tìm đường dẫn tuyệt đối tới ffmpeg
    ffmpeg_bin = shutil.which("ffmpeg")
    if not ffmpeg_bin:
        possible_ffmpeg = os.path.join(os.path.dirname(sys.executable), "ffmpeg.exe")
        if os.path.exists(possible_ffmpeg):
            ffmpeg_bin = possible_ffmpeg
        else:
            ffmpeg_bin = "ffmpeg"

    try:
        comtypes.CoInitialize()
    except Exception:
        pass

    out_dir = _temp_dir()
    pptx_abs = str(Path(pptx_path).resolve())
    temp_video_path = os.path.join(out_dir, "full_presentation.mp4")

    report(5, "Đang kết nối PowerPoint...")
    ppt_app = None
    try:
        ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
        # ppt_app.Visible = 1
        prs = ppt_app.Presentations.Open(pptx_abs, ReadOnly=True, WithWindow=False)
        slide_count = prs.Slides.Count
        report(10, f"Đang yêu cầu PowerPoint xuất {slide_count} slide thành video (quá trình này mất khoảng 20-40s)...")

        # 5.0 giây cho mỗi slide làm mốc thời lượng cố định để cắt
        SLIDE_DUR = 5.0

        # CreateVideo(FileName, UseSlidesAndTimings, DefaultSlideDuration, Resolution, FrameRate, Quality)
        # UseSlidesAndTimings=False để ép tất cả slide đều dài SLIDE_DUR giây
        prs.CreateVideo(temp_video_path, False, int(SLIDE_DUR), 720, 30, 85)

        # Chờ xuất video hoàn thành
        start_wait = time.time()
        while True:
            if time.time() - start_wait > 300:  # Timeout 5 phút
                raise RuntimeError("PowerPoint xuất video bị quá thời gian chờ (5 phút).")

            status = prs.CreateVideoStatus
            if status == 3:  # Done
                break
            elif status == 4:  # Failed
                raise RuntimeError("PowerPoint COM báo lỗi xuất video.")

            elapsed = int(time.time() - start_wait)
            report(15 + min(50, elapsed), f"Đang xuất video slide ({elapsed}s)...")
            time.sleep(1.0)

        report(70, f"Đang cắt video lớn thành các slide riêng biệt…")
        
        slides: List[SlideInfo] = []
        titles = _extract_pptx_titles(pptx_path)

        for i in range(1, slide_count + 1):
            start_time = (i - 1) * SLIDE_DUR
            slide_video_path = os.path.join(out_dir, f"slide_{i:04d}.mp4")
            img_path = os.path.join(out_dir, f"slide_{i:04d}.png")

            # Cắt phân đoạn video
            subprocess.run([
                ffmpeg_bin, "-y", "-ss", f"{start_time:.2f}", "-t", f"{SLIDE_DUR:.2f}",
                "-i", temp_video_path,
                "-c:v", "libx264", "-pix_fmt", "yuv420p", "-an",
                slide_video_path
            ], capture_output=True)

            # Xuất 1 frame làm ảnh đại diện/thumbnail
            subprocess.run([
                ffmpeg_bin, "-y", "-ss", "00:00:01.00", "-i", slide_video_path,
                "-vframes", "1", "-vf", "scale=320:-1", img_path
            ], capture_output=True)

            title = titles[i - 1] if (i - 1) < len(titles) else ""
            slides.append(SlideInfo(
                index=i - 1,
                image_path=img_path,
                video_path=slide_video_path,
                title=title
            ))
            
            pct = 70 + int(25 * i / slide_count)
            report(pct, f"Slide {i}/{slide_count}…")

        prs.Close()
    finally:
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass
        try:
            comtypes.CoUninitialize()
        except Exception:
            pass

    # Xóa file video gốc tạm thời
    if os.path.exists(temp_video_path):
        try:
            os.remove(temp_video_path)
        except Exception:
            pass

    report(100, f"Đã chuẩn bị {len(slides)} slide dạng video thành công.")
    return slides


def _pptx_via_com_export_slides(
    pptx_path: str,
    report: Callable[[int, str], None],
) -> List[SlideInfo]:
    """Dùng PowerPoint COM Export để kết xuất ảnh PNG gốc và trích xuất ảnh GIF nhúng bằng python-pptx."""
    import comtypes.client  # type: ignore
    import comtypes
    import time

    try:
        comtypes.CoInitialize()
    except Exception:
        pass

    out_dir = _temp_dir()
    pptx_abs = str(Path(pptx_path).resolve())

    report(5, "Kết nối PowerPoint...")
    ppt_app = None
    try:
        ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
        ppt_app.Visible = 1  # phải visible để render đúng
        prs = ppt_app.Presentations.Open(pptx_abs, ReadOnly=True, WithWindow=False)
        slide_count = prs.Slides.Count
        report(10, f"Đang kết xuất {slide_count} slide ảnh gốc sắc nét...")

        # Trích xuất thông tin ảnh GIF bằng python-pptx
        slide_gifs_map = {}
        try:
            from pptx import Presentation
            pptx_prs = Presentation(pptx_path)
            prs_w = pptx_prs.slide_width
            prs_h = pptx_prs.slide_height

            for idx, pptx_slide in enumerate(pptx_prs.slides):
                gifs = []
                for shape in pptx_slide.shapes:
                    shape_type = getattr(shape, "shape_type", 0)
                    # 13 là MSO_SHAPE_TYPE.PICTURE
                    if shape_type == 13:
                        try:
                            img = shape.image
                            if img.ext.lower() == "gif":
                                gif_filename = f"slide_{idx:04d}_gif_{len(gifs)}.gif"
                                gif_path = os.path.join(out_dir, gif_filename)
                                with open(gif_path, "wb") as f:
                                    f.write(img.blob)

                                x_rel = shape.left / prs_w
                                y_rel = shape.top / prs_h
                                w_rel = shape.width / prs_w
                                h_rel = shape.height / prs_h

                                gifs.append({
                                    "gif_path": gif_path,
                                    "x_rel": x_rel,
                                    "y_rel": y_rel,
                                    "w_rel": w_rel,
                                    "h_rel": h_rel
                                })
                        except Exception as ex:
                            print(f"Lỗi trích xuất GIF: {ex}")
                if gifs:
                    slide_gifs_map[idx] = gifs
        except Exception as ex:
            print(f"Lỗi phân tích python-pptx: {ex}")

        slides: List[SlideInfo] = []
        titles = _extract_pptx_titles(pptx_path)

        for i in range(1, slide_count + 1):
            img_path = os.path.join(out_dir, f"slide_{i:04d}.png")
            slide = prs.Slides(i)
            # slide.Export xuất ảnh gốc cực kỳ sắc nét của PowerPoint
            slide.Export(img_path, "PNG")
            title = titles[i - 1] if (i - 1) < len(titles) else ""
            
            gifs = slide_gifs_map.get(i - 1, [])
            slides.append(SlideInfo(index=i - 1, image_path=img_path, title=title, gifs=gifs))
            
            pct = 10 + int(85 * i / slide_count)
            report(pct, f"Slide {i}/{slide_count}...")

        prs.Close()
    finally:
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass
        try:
            comtypes.CoUninitialize()
        except Exception:
            pass

    report(100, f"Đã kết xuất {len(slides)} slide sắc nét.")
    return slides


def _pptx_via_com_to_pdf(
    pptx_path: str,
    report: Callable[[int, str], None],
) -> List[SlideInfo]:
    """Convert PPTX → PDF qua COM, rồi render PDF bằng PyMuPDF."""
    import comtypes.client  # type: ignore
    import comtypes

    try:
        comtypes.CoInitialize()
    except Exception:
        pass

    pdf_path = os.path.join(_temp_dir(), "slides.pdf")
    pptx_abs = str(Path(pptx_path).resolve())
    pdf_abs  = str(Path(pdf_path).resolve())

    report(5, "Kết nối PowerPoint để convert PDF…")
    ppt_app = None
    try:
        ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
        ppt_app.Visible = 1
        prs = ppt_app.Presentations.Open(pptx_abs, ReadOnly=True, WithWindow=False)
        # ppSaveAsPDF = 32
        prs.SaveAs(pdf_abs, 32)
        prs.Close()
    finally:
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass
        try:
            comtypes.CoUninitialize()
        except Exception:
            pass

    report(40, "Convert xong, đang render PDF…")
    return _render_pdf(pdf_abs, report, start_pct=40, titles=_extract_pptx_titles(pptx_path))


def _pptx_fallback_text(
    pptx_path: str,
    report: Callable[[int, str], None],
) -> List[SlideInfo]:
    """
    Fallback cuối: tạo ảnh placeholder màu tối kèm text từ python-pptx.
    Dùng Pillow để vẽ text lên ảnh.
    """
    try:
        from pptx import Presentation  # type: ignore
        from PIL import Image, ImageDraw, ImageFont  # type: ignore
    except ImportError:
        raise SlideProcessorError(
            "Không thể render PPTX: thiếu thư viện python-pptx và/hoặc Pillow."
        )

    out_dir = _temp_dir()
    prs = Presentation(pptx_path)
    slides: List[SlideInfo] = []

    for i, slide in enumerate(prs.slides, start=1):
        # Lấy text từ slide
        lines: List[str] = []
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if not title and slide.shapes.title and shape == slide.shapes.title:
                        title = text
                    lines.append(text)

        # Vẽ placeholder
        img = Image.new("RGB", (1280, 720), color=(22, 27, 34))
        draw = ImageDraw.Draw(img)

        # Slide number badge
        draw.rectangle([30, 30, 110, 70], fill=(109, 40, 217))
        draw.text((50, 38), f"#{i}", fill=(255, 255, 255))

        # Title
        y = 100
        for j, line in enumerate(lines[:10]):
            draw.text((60, y), line[:80], fill=(230, 237, 243))
            y += 40
            if y > 650:
                break

        img_path = os.path.join(out_dir, f"slide_{i:04d}.png")
        img.save(img_path, "PNG")
        slides.append(SlideInfo(index=i - 1, image_path=img_path, title=title))
        report(int(100 * i / len(prs.slides)), f"Slide {i}…")

    return slides


# ═══════════════════════════════════════════════════════════════════════════
#  PDF → THUMBNAILS  (PyMuPDF)
# ═══════════════════════════════════════════════════════════════════════════

def load_pdf(
    pdf_path: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
) -> List[SlideInfo]:
    """Render từng trang PDF thành PNG thumbnail qua PyMuPDF."""

    def _report(pct: int, msg: str):
        if progress_cb:
            progress_cb(pct, msg)

    _report(0, "Đang mở file PDF…")
    return _render_pdf(pdf_path, _report)


def _render_pdf(
    pdf_path: str,
    report: Callable[[int, str], None],
    start_pct: int = 0,
    titles: Optional[List[str]] = None,
) -> List[SlideInfo]:
    """Shared helper: render PDF → PNG list dùng PyMuPDF."""
    try:
        import fitz  # PyMuPDF  # type: ignore
    except ImportError:
        raise SlideProcessorError(
            "Thiếu thư viện PyMuPDF. Chạy: pip install PyMuPDF"
        )

    out_dir = _temp_dir()
    doc = fitz.open(pdf_path)
    page_count = len(doc)
    slides: List[SlideInfo] = []

    for i, page in enumerate(doc, start=1):
        mat = fitz.Matrix(1280 / page.rect.width, 720 / page.rect.height)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_path = os.path.join(out_dir, f"slide_{i:04d}.png")
        pix.save(img_path)

        title = (titles[i - 1] if titles and (i - 1) < len(titles) else "")
        slides.append(SlideInfo(index=i - 1, image_path=img_path, title=title))

        pct = start_pct + int((100 - start_pct) * i / page_count)
        report(pct, f"Trang {i}/{page_count}…")

    doc.close()
    report(100, f"Đã render {len(slides)} trang.")
    return slides


# ═══════════════════════════════════════════════════════════════════════════
#  PROJECT SAVE / LOAD
# ═══════════════════════════════════════════════════════════════════════════

def mapping_to_dict(slides: List[SlideInfo]) -> dict:
    """Chuyển danh sách slide assignments sang dict để lưu JSON."""
    return {
        "version": 1,
        "assignments": [
            {
                "slide_index":   s.index,
                "slide_title":   s.title,
                "assigned_pos":  s.assigned_pos,
                "assigned_text": s.assigned_text,
            }
            for s in slides
        ],
    }


def apply_mapping_from_dict(slides: List[SlideInfo], data: dict) -> None:
    """Áp dụng mapping đã lưu vào danh sách slide hiện tại."""
    if data.get("version") != 1:
        return
    index_map = {s.index: s for s in slides}
    for item in data.get("assignments", []):
        idx = item.get("slide_index", -1)
        if idx in index_map:
            index_map[idx].assigned_pos  = item.get("assigned_pos", -1)
            index_map[idx].assigned_text = item.get("assigned_text", "")
